from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Define Theme Colors
    BLUE_NEXUS = RGBColor(59, 130, 246)
    DARK_NEXUS = RGBColor(2, 6, 23)

    # Helper function to style titles
    def style_title(slide, text):
        title = slide.shapes.title
        title.text = text
        for paragraph in title.text_frame.paragraphs:
            paragraph.font.size = Pt(36)
            paragraph.font.bold = True
            paragraph.font.color.rgb = BLUE_NEXUS
            paragraph.alignment = PP_ALIGN.LEFT

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Clinical Nexus"
    subtitle.text = "Agentic AI Operating System for Precision Medicine\nPresented by: Deepak Kumar Yadav"
    
    # Slide 2: The Core Problem
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    style_title(slide, "The Problem: Institutional Inefficiency")
    content = slide.placeholders[1]
    content.text = ("- Medical documentation is a $30B manual burden\n"
                  "- Clinical data is fragmented across PDFs and disparate case files\n"
                  "- Diagnostic delays due to information overload and high provider burnout")

    # Slide 3: Our Solution
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    style_title(slide, "Nexus: The Agentic Solution")
    content = slide.placeholders[1]
    content.text = ("- A Multi-Modal AI Operating System for healthcare\n"
                  "- Real-time Disease Prediction via robust ML Engines\n"
                  "- Context-Aware Intelligence Nodes: Scribe, Consult, and Prescribe")

    # Slide 4: Key Feature - Neural OCR Pipeline
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    style_title(slide, "Key Innovation: Neural OCR Hub")
    content = slide.placeholders[1]
    content.text = ("- OpenCV Pre-processing for noise reduction in lab reports\n"
                  "- Tesseract-powered extraction of sub-second biometric values\n"
                  "- Direct sync to Diagnostic Models for immediate analysis")

    # Slide 5: Key Feature - Chief Clinical Officer (CCO) Mode
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    style_title(slide, "Phase 9: Chief Clinical Intelligence")
    content = slide.placeholders[1]
    content.text = ("- Expert Role-Prompting using Claude 3.5 Sonnet/Haiku\n"
                  "- Recursive medical knowledge graph (PubMed, Lancet, JAMA)\n"
                  "- Interactive 'Doctor-Patient' Consultation interface for real-time dialogue")

    # Slide 6: Key Feature - Precision AI Scribe
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    style_title(slide, "Precision AI Scribe & Prescription Suite")
    content = slide.placeholders[1]
    content.text = ("- Automated SOAP note generation (Subjective, Objective, Assessment, Plan)\n"
                  "- Data-driven medication suggestions based on diagnostic results\n"
                  "- Professional PDF report generation for physician-in-the-loop oversight")

    # Slide 7: Technical Stack & Scalability
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    style_title(slide, "Enterprise Architecture (Nexus Grid)")
    content = slide.placeholders[1]
    content.text = ("- Backend: FastAPI (Python 3.12) / SQLAlchemy\n"
                  "- Frontend: React 19 / Vite 6 / Tailwind CSS v4 / Framer Motion\n"
                  "- Scalability: Modular Nodes & Neural Security for PHI HIPAA-alignment")

    # Slide 8: Closing
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    style_title(slide, "Clinical Nexus: Precision in Action")
    content = slide.placeholders[1]
    content.text = ("- Enhancing clinical outcomes via Agentic AI\n"
                  "- Status: PROTOTYPE SYNCED & LIVE\n"
                  "- Presented by: Deepak Kumar Yadav")

    # Save the file
    file_name = 'Clinical_Nexus_Presentation.pptx'
    prs.save(file_name)
    print(f"✅ Slide Deck Generated: {file_name}")

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("❌ Error: python-pptx not found. Run 'python -m pip install python-pptx'")
